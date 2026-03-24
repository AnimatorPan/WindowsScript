"""
Office 文档预览
"""
from PyQt6.QtWidgets import QVBoxLayout, QTextEdit, QLabel
from PyQt6.QtCore import Qt
from pathlib import Path
import logging

from .preview_base import PreviewWidget

logger = logging.getLogger(__name__)


class OfficePreview(PreviewWidget):
    """Office 文档预览"""

    def init_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)

        self.text_edit = QTextEdit()
        self.text_edit.setReadOnly(True)
        self.text_edit.setAcceptRichText(True)
        layout.addWidget(self.text_edit)

    def do_preview(self, file_path: str):
        ext = Path(file_path).suffix.lower()

        if ext in {'.doc', '.docx'}:
            self.preview_word(file_path)
        elif ext in {'.xls', '.xlsx'}:
            self.preview_excel(file_path)
        elif ext in {'.ppt', '.pptx'}:
            self.preview_powerpoint(file_path)
        elif ext == '.rtf':
            self.preview_rtf(file_path)
        elif ext in {'.odt', '.ods', '.odp'}:
            self.preview_openoffice(file_path)
        else:
            self.show_error(f"不支持的文件格式: {ext}")

    def preview_word(self, file_path: str):
        try:
            from docx import Document

            doc = Document(file_path)
            content = []

            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text)

            for table in doc.tables:
                content.append("\n[表格]")
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    content.append(row_text)

            text = "\n".join(content)
            if len(text) > 50000:
                text = text[:50000] + "\n\n... (内容过长，仅显示部分)"

            self.text_edit.setPlainText(text)

        except ImportError:
            self.text_edit.setHtml(
                "<div style='padding: 20px; text-align: center;'>"
                "<h3>Word 文档预览</h3>"
                "<p>需要安装 python-docx 库</p>"
                "<p style='color: #888;'>请运行: pip install python-docx</p>"
                "</div>"
            )
        except Exception as e:
            logger.error(f"预览 Word 文档失败: {e}")
            self.show_error(f"预览失败: {e}")

    def preview_excel(self, file_path: str):
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, read_only=True, data_only=True)
            content = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                content.append(f"\n=== {sheet_name} ===\n")

                row_count = 0
                for row in sheet.iter_rows(max_row=100, values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                        content.append(row_text)
                        row_count += 1

                if row_count >= 100:
                    content.append("... (仅显示前 100 行)")

            wb.close()
            text = "\n".join(content)
            self.text_edit.setPlainText(text)

        except ImportError:
            self.text_edit.setHtml(
                "<div style='padding: 20px; text-align: center;'>"
                "<h3>Excel 文档预览</h3>"
                "<p>需要安装 openpyxl 库</p>"
                "<p style='color: #888;'>请运行: pip install openpyxl</p>"
                "</div>"
            )
        except Exception as e:
            logger.error(f"预览 Excel 文档失败: {e}")
            self.show_error(f"预览失败: {e}")

    def preview_powerpoint(self, file_path: str):
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            content = ["=== PowerPoint 文档预览 ===\n"]

            for i, slide in enumerate(prs.slides, 1):
                content.append(f"\n--- 幻灯片 {i} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(shape.text)

            text = "\n".join(content)
            self.text_edit.setPlainText(text)

        except ImportError:
            self.text_edit.setHtml(
                "<div style='padding: 20px; text-align: center;'>"
                "<h3>PowerPoint 文档预览</h3>"
                "<p>需要安装 python-pptx 库</p>"
                "<p style='color: #888;'>请运行: pip install python-pptx</p>"
                "</div>"
            )
        except Exception as e:
            logger.error(f"预览 PowerPoint 文档失败: {e}")
            self.show_error(f"预览失败: {e}")

    def preview_rtf(self, file_path: str):
        try:
            from striprtf.striprtf import rtf_to_text

            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()

            text = rtf_to_text(rtf_content)
            if len(text) > 50000:
                text = text[:50000] + "\n\n... (内容过长，仅显示部分)"

            self.text_edit.setPlainText(text)

        except ImportError:
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                self.text_edit.setPlainText(f"[RTF 原始内容]\n\n{content[:5000]}")
            except Exception as e:
                self.show_error(f"预览失败: {e}")
        except Exception as e:
            logger.error(f"预览 RTF 文档失败: {e}")
            self.show_error(f"预览失败: {e}")

    def preview_openoffice(self, file_path: str):
        ext = Path(file_path).suffix.lower()

        if ext == '.odt':
            try:
                from odf.opendocument import load
                from odf import text

                doc = load(file_path)
                content = []

                for element in doc.getElementsByType(text.P):
                    if element.firstChild:
                        content.append(str(element.firstChild))

                text_content = "\n".join(content)
                self.text_edit.setPlainText(text_content)

            except ImportError:
                self.text_edit.setHtml(
                    "<div style='padding: 20px; text-align: center;'>"
                    "<h3>OpenDocument 文档预览</h3>"
                    "<p>需要安装 odfpy 库</p>"
                    "<p style='color: #888;'>请运行: pip install odfpy</p>"
                    "</div>"
                )
            except Exception as e:
                logger.error(f"预览 OpenDocument 失败: {e}")
                self.show_error(f"预览失败: {e}")
        else:
            self.text_edit.setPlainText(
                f"OpenDocument 格式 ({ext}) 预览\n\n"
                f"建议使用专业工具打开此文件"
            )

    def clear(self):
        self.text_edit.clear()
        self.current_file = None
